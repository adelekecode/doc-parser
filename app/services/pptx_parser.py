import os
from pptx import Presentation
from app.utils.exceptions import ParsingError







class PPTXParser:
    def parse(self, file_path):
        """
        Parse PPTX document and extract relevant information
        
        Args:
            file_path (str): Path to the PPTX file
            
        Returns:
            dict: Extracted data from the PPTX
            
        Raises:
            ParsingError: If there's an error during parsing
        """
        if not os.path.exists(file_path):
            raise ParsingError(f"File not found: {file_path}")
        
        try:
            prs = Presentation(file_path)
            
            # Extract document metadata
            metadata = {
                'author': prs.core_properties.author,
                'title': prs.core_properties.title,
                'subject': prs.core_properties.subject,
                'created': prs.core_properties.created.isoformat() if prs.core_properties.created else None,
                'modified': prs.core_properties.modified.isoformat() if prs.core_properties.modified else None,
                'category': prs.core_properties.category,
                'comments': prs.core_properties.comments,
                'keywords': prs.core_properties.keywords,
                'last_modified_by': prs.core_properties.last_modified_by,
                'revision': prs.core_properties.revision
            }
            
            # Remove None values
            metadata = {k: v for k, v in metadata.items() if v is not None}
            
            # Extract content from each slide
            slides = []
            for i, slide in enumerate(prs.slides):
                # Extract title
                title = None
                for shape in slide.shapes:
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        if shape.text.strip():
                            title = shape.text.strip()
                            break
                
                # Extract text content
                content = []
                for shape in slide.shapes:
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        if shape.text.strip():
                            content.append(shape.text.strip())
                
                # Check for images and charts
                has_images = False
                has_charts = False
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        has_images = True
                    elif shape.shape_type == 3:  # MSO_SHAPE_TYPE.CHART
                        has_charts = True
                
                slides.append({
                    'slide_number': i + 1,
                    'title': title,
                    'content': '\n'.join(content),
                    'has_images': has_images,
                    'has_charts': has_charts
                })
            
            return {
                'total_slides': len(prs.slides),
                'slides': slides,
                'metadata': metadata
            }
            
        except Exception as e:
            raise ParsingError(f"Error parsing PPTX: {str(e)}")
